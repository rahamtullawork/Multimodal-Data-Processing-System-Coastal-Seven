# file_parser.py
import fitz  # PyMuPDF for PDFs
from pdfminer.high_level import extract_text
import docx
from pptx import Presentation
import os


def extract_text_from_file(file_path):
    """
    Extracts text from supported file types:
    PDF, DOCX, PPTX, TXT, MD
    Returns extracted text as a string.
    """
    ext = os.path.splitext(file_path)[1].lower()
    text = ""

    try:
        if ext == ".pdf":
            text = extract_text_from_pdf(file_path)
        elif ext == ".docx":
            text = extract_text_from_docx(file_path)
        elif ext == ".pptx":
            text = extract_text_from_pptx(file_path)
        elif ext in [".txt", ".md"]:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
        else:
            text = ""
    except Exception as e:
        print(f"[ERROR] Failed to extract from {file_path}: {e}")

    return text.strip()


def extract_text_from_pdf(file_path):
    """Extract text from PDF using pdfminer"""
    return extract_text(file_path)


def extract_text_from_docx(file_path):
    """Extract text from DOCX"""
    doc = docx.Document(file_path)
    return "\n".join([p.text for p in doc.paragraphs])


def extract_text_from_pptx(file_path):
    """Extract text from PowerPoint"""
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text
